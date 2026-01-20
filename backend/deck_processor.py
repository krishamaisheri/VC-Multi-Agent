import base64
import io
import json
import logging
import os
import tempfile
from typing import List, Dict, Optional
from uuid import uuid4

import pdfplumber
from PIL import Image
from pptx import Presentation

logger = logging.getLogger(__name__)


def _describe_images_with_vlm(image_buffers: List[bytes], call_vision_fn) -> List[str]:
    descriptions = []
    for idx, buf in enumerate(image_buffers):
        try:
            b64 = base64.b64encode(buf).decode("utf-8")
            prompt = "Describe this slide image briefly for a VC pitch."
            desc = call_vision_fn(b64, prompt)
            descriptions.append(desc or "")
        except Exception as e:  # noqa: BLE001
            logger.error(f"VLM description failed for image {idx}: {e}")
            descriptions.append("")
    return descriptions


def _ppt_extract_images(slide) -> List[bytes]:
    buffers = []
    for shape in slide.shapes:
        if hasattr(shape, "image"):
            try:
                image = shape.image
                stream = io.BytesIO(image.blob)
                buffers.append(stream.read())
            except Exception as e:  # noqa: BLE001
                logger.warning(f"Could not extract image from shape: {e}")
    return buffers


def extract_from_pptx(file_bytes: bytes, call_vision_fn) -> List[Dict]:
    prs = Presentation(io.BytesIO(file_bytes))
    pages = []
    for i, slide in enumerate(prs.slides, start=1):
        text_runs = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
        images = _ppt_extract_images(slide)
        image_desc = _describe_images_with_vlm(images, call_vision_fn) if images else []
        content_parts = []
        if text_runs:
            content_parts.append("\n".join(filter(None, text_runs)))
        if image_desc:
            joined_desc = "\n".join([f"Image {idx+1}: {d}" for idx, d in enumerate(image_desc) if d])
            if joined_desc:
                content_parts.append("[Image Descriptions]\n" + joined_desc)
        pages.append({
            "page_number": i,
            "content": "\n\n".join(content_parts) if content_parts else "",
            "has_images": bool(images),
        })
    return pages


def extract_from_pdf(file_bytes: bytes, call_vision_fn) -> List[Dict]:
    pages = []
    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            text = page.extract_text() or ""
            image_desc = []
            # pdfplumber gives image metadata; attempt to render for vision
            try:
                for img in page.images:
                    try:
                        x0, y0, x1, y1 = img.get("x0"), img.get("y0"), img.get("x1"), img.get("y1")
                        if None in (x0, y0, x1, y1):
                            continue
                        cropped = page.crop((x0, y0, x1, y1)).to_image(resolution=200)
                        buf = io.BytesIO()
                        cropped.save(buf, format="PNG")
                        image_desc.extend(_describe_images_with_vlm([buf.getvalue()], call_vision_fn))
                    except Exception as e:  # noqa: BLE001
                        logger.warning(f"Failed to process pdf image on page {i}: {e}")
            except Exception as e:  # noqa: BLE001
                logger.warning(f"Image extraction not available for pdf page {i}: {e}")
            content_parts = []
            if text:
                content_parts.append(text)
            if image_desc:
                joined_desc = "\n".join([f"Image {idx+1}: {d}" for idx, d in enumerate(image_desc) if d])
                if joined_desc:
                    content_parts.append("[Image Descriptions]\n" + joined_desc)
            pages.append({
                "page_number": i,
                "content": "\n\n".join(content_parts) if content_parts else "",
                "has_images": bool(image_desc),
            })
    return pages


def process_pitch_file(file_name: Optional[str], file_base64: Optional[str], call_vision_fn) -> List[Dict]:
    if not file_name or not file_base64:
        return []
    try:
        file_bytes = base64.b64decode(file_base64)
    except Exception as e:  # noqa: BLE001
        logger.error(f"Invalid base64 for pitch file: {e}")
        return []

    lower = file_name.lower()
    if lower.endswith((".pptx", ".ppt")):
        try:
            return extract_from_pptx(file_bytes, call_vision_fn)
        except Exception as e:  # noqa: BLE001
            logger.error(f"PPTX extraction failed: {e}")
            return []
    if lower.endswith(".pdf"):
        try:
            return extract_from_pdf(file_bytes, call_vision_fn)
        except Exception as e:  # noqa: BLE001
            logger.error(f"PDF extraction failed: {e}")
            return []
    logger.warning(f"Unsupported pitch file type: {file_name}")
    return []


def build_page_documents(pages: List[Dict], metadata_base: Dict) -> List[Dict]:
    docs = []
    for page in pages:
        content = page.get("content", "").strip()
        if not content:
            continue
        doc = {
            "id": str(uuid4()),
            "text": content,
            "metadata": {
                **metadata_base,
                "page_number": page.get("page_number"),
                "has_images": page.get("has_images", False),
            }
        }
        docs.append(doc)
    return docs
